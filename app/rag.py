import os
import shutil
import time
import requests
from typing import List, Dict, Optional
from langchain_community.document_loaders import PyPDFLoader, WebBaseLoader
# ... (rest of imports)
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_core.embeddings import Embeddings
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_core.documents import Document
import docx
from pptx import Presentation

try:
    from .config import config
except ImportError:
    from config import config

# Global dictionary to store FAISS index objects in memory keyed by course_id_division
faiss_indexes: Dict[str, FAISS] = {}

# Initialize Local Embeddings
embeddings = HuggingFaceEmbeddings(
    model_name='sentence-transformers/all-MiniLM-L6-v2'
)

# Text Splitter
text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)

def get_index_path(course_id: str, division: str):
    """Returns the persistent folder path for this course and division's FAISS index."""
    os.makedirs(config.FAISS_DB_PATH, exist_ok=True)
    return config.FAISS_DB_PATH / f"{course_id}_{division}"

def get_vector_store(course_id: str, division: str):
    """Returns FAISS index from memory Cache or loads from Disk."""
    key = f"{course_id}_{division}"
    
    # Check memory cache first
    if key in faiss_indexes:
        return faiss_indexes[key]
        
    # Check disk
    index_path = get_index_path(course_id, division)
    if os.path.exists(os.path.join(index_path, "index.faiss")):
        try:
            vector_store = FAISS.load_local(str(index_path), embeddings, allow_dangerous_deserialization=True)
            faiss_indexes[key] = vector_store
            return vector_store
        except Exception as e:
            print(f"Failed to load FAISS index from disk: {e}")
            
    return None

def load_document(file_path: str, filename: str) -> List[Document]:
    ext = os.path.splitext(filename)[1].lower()
    if ext == '.pdf':
        loader = PyPDFLoader(file_path)
        return loader.load()
    elif ext in ['.pptx', '.ppt']:
        prs = Presentation(file_path)
        docs = []
        for i, slide in enumerate(prs.slides):
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text.append(paragraph.text)
            if text:
                content = "\n".join(text)
                docs.append(Document(page_content=content, metadata={"source": filename, "slide": i + 1, "type": "pptx"}))
        return docs
    elif ext in ['.docx', '.doc']:
        doc = docx.Document(file_path)
        text = []
        for para in doc.paragraphs:
            if para.text.strip():
                text.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text.append(cell.text)
        content = "\n".join(text)
        if not content:
            return []
        return [Document(page_content=content, metadata={"source": filename, "type": "docx"})]
    else:
        raise ValueError(f"Unsupported file type: {ext}")

def process_document(file_path: str, course_id: str, division: str):
    """Loads a Document, chunks it, and adds to FAISS index for course+division."""
    import datetime
    filename = os.path.basename(file_path)
    docs = load_document(file_path, filename)
    if not docs:
        return 0
    chunks = text_splitter.split_documents(docs)
    
    today_str = datetime.date.today().strftime("%Y-%m-%d")
    
    # Add metadata
    for chunk in chunks:
        # Don't overwrite type if docx or pptx
        meta = {"course_id": course_id, "division": division, "source": filename, "session_date": today_str}
        if "type" not in chunk.metadata:
            meta["type"] = "pdf"
        chunk.metadata.update(meta)
    
    vector_store = get_vector_store(course_id, division)
    if vector_store:
        vector_store.add_documents(chunks)
    else:
        vector_store = FAISS.from_documents(chunks, embeddings)
        
    # Save to disk
    vector_store.save_local(str(get_index_path(course_id, division)))
    
    # Update cache
    key = f"{course_id}_{division}"
    faiss_indexes[key] = vector_store
    
    return len(chunks)

def process_url(url: str, course_id: str, division: str):
    """Scrapes a URL, chunks it, and adds to FAISS index for course+division."""
    # Use a more standard User-Agent to avoid being blocked
    loader = WebBaseLoader(
        url,
        header_template={
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
    )
    log_path = os.path.join(os.path.dirname(__file__), "..", "rag_debug.log")
    with open(log_path, "a", encoding="utf-8") as f:
        f.write(f"\n--- {time.ctime()} --- URL: {url}\n")
        try:
            docs = loader.load()
            if not docs:
                f.write("DEBUG: No content scraped.\n")
                return 0
            
            content_snippet = docs[0].page_content[:500].replace('\n', ' ')
            f.write(f"DEBUG: Scraped {len(docs)} documents. Snippet: {content_snippet}\n")
            
            chunks = text_splitter.split_documents(docs)
            f.write(f"DEBUG: Split into {len(chunks)} chunks.\n")
        except Exception as e:
            f.write(f"DEBUG: Scraping Error: {str(e)}\n")
            return 0
    
    import datetime
    today_str = datetime.date.today().strftime("%Y-%m-%d")
    
    # Add metadata
    for chunk in chunks:
        chunk.metadata.update({"course_id": course_id, "division": division, "source": url, "session_date": today_str})
    
    vector_store = get_vector_store(course_id, division)
    if vector_store:
        vector_store.add_documents(chunks)
    else:
        vector_store = FAISS.from_documents(chunks, embeddings)
        
    # Save to disk
    vector_store.save_local(str(get_index_path(course_id, division)))
    
    # Update cache
    key = f"{course_id}_{division}"
    faiss_indexes[key] = vector_store
    
    return len(chunks)

def get_context(course_id: str, division: str, query: Optional[str] = None, n_results: int = 10):
    """Retrieves relevant chunks from the specific FAISS index for course+division."""
    vector_store = get_vector_store(course_id, division)
    if not vector_store:
        return ""
    
    if query:
        k = min(n_results, vector_store.index.ntotal)
        results = vector_store.similarity_search(query, k=k) if k > 0 else []
    else:
        # Return most recent chunks
        k = min(n_results, vector_store.index.ntotal)
        results = vector_store.similarity_search("", k=k) if k > 0 else []
        
    return "\n\n".join([doc.page_content for doc in results])

def get_context_filtered(course_id: str, division: str, query: str = "", session_date: Optional[str] = None, max_chars: int = 20000) -> str:
    """Retrieves chunks proportionally from the docstore, optionally filtered by session_date."""
    vector_store = get_vector_store(course_id, division)
    if not vector_store:
        return ""
        
    k = min(40, vector_store.index.ntotal)
    filter_dict = {"session_date": session_date} if session_date else {}
    results = vector_store.similarity_search(query, k=k, filter=filter_dict) if k > 0 else []

    from collections import defaultdict
    by_source = defaultdict(list)
    for doc in results:
        source = doc.metadata.get('source', 'unknown')
        by_source[source].append(doc.page_content)

    if not by_source:
        return ""

    source_sizes = {
        source: sum(len(chunk) for chunk in chunks)
        for source, chunks in by_source.items()
    }
    total_size = sum(source_sizes.values())
    if total_size == 0:
        return ""

    num_sources = len(by_source)
    min_budget = min(int(max_chars * 0.10), max_chars // (num_sources * 2))
    remaining_budget = max_chars - (min_budget * num_sources)

    context_parts = []
    for source, chunks in by_source.items():
        proportional_share = (source_sizes[source] / total_size) * remaining_budget
        budget = int(min_budget + proportional_share)
        source_text = " ".join(chunks)[:budget]
        context_parts.append(f"[From {source}]\n{source_text}")
        print(f"DEBUG: '{source}' filter={session_date} size={source_sizes[source]} budget={budget} chars")

    return "\n\n".join(context_parts)

def get_all_context(course_id: str, division: str, query: str = "", max_chars: int = 20000) -> str:
    """Retrieves all chunks proportionally from the docstore for comprehensive tasks like summarization."""
    vector_store = get_vector_store(course_id, division)
    if not vector_store:
        return ""
    
    # Get a large pool of relevant chunks
    k = min(40, vector_store.index.ntotal)
    results = vector_store.similarity_search(query, k=k) if k > 0 else []

    # Group chunks by source and calculate total size per source
    from collections import defaultdict
    by_source = defaultdict(list)
    for doc in results:
        source = doc.metadata.get('source', 'unknown')
        by_source[source].append(doc.page_content)

    if not by_source:
        return ""

    # Calculate total chars per source
    source_sizes = {
        source: sum(len(chunk) for chunk in chunks)
        for source, chunks in by_source.items()
    }
    total_size = sum(source_sizes.values())
    if total_size == 0:
        return ""

    # Guarantee minimum 10% of budget per source regardless of size
    num_sources = len(by_source)
    min_budget = min(int(max_chars * 0.10), max_chars // (num_sources * 2))
    remaining_budget = max_chars - (min_budget * num_sources)

    # Allocate remaining budget proportionally by document size
    context_parts = []
    for source, chunks in by_source.items():
        proportional_share = (source_sizes[source] / total_size) * remaining_budget
        budget = int(min_budget + proportional_share)
        source_text = " ".join(chunks)[:budget]
        context_parts.append(f"[From {source}]\n{source_text}")
        print(f"DEBUG: '{source}' size={source_sizes[source]} budget={budget} chars")

    return "\n\n".join(context_parts)

def remove_session_date(course_id: str, division: str, session_date: str):
    """Removes all chunks associated with a specific session date from the FAISS index."""
    vector_store = get_vector_store(course_id, division)
    if not vector_store:
        return 0
    
    all_docs = list(vector_store.docstore._dict.values())
    remaining_docs = [doc for doc in all_docs if doc.metadata.get('session_date') != session_date]
    
    if len(remaining_docs) == len(all_docs):
        return 0
    
    clear_session_material(course_id, division)
    
    if remaining_docs:
        new_vs = FAISS.from_documents(remaining_docs, embeddings)
        key = f"{course_id}_{division}"
        faiss_indexes[key] = new_vs
        new_vs.save_local(str(get_index_path(course_id, division)))
    
    return len(all_docs) - len(remaining_docs)

def remove_source(course_id: str, division: str, source_name: str):
    """Removes a specific source (file or URL) from the FAISS index."""
    vector_store = get_vector_store(course_id, division)
    if not vector_store:
        return 0
    
    # Extract all documents from the vector store
    # This is a bit of a hack but FAISS doesn't expose doc list easily
    # We use the internal docstore
    all_docs = list(vector_store.docstore._dict.values())
    
    # Filter out docs from the target source
    remaining_docs = [doc for doc in all_docs if doc.metadata.get('source') != source_name]
    
    if len(remaining_docs) == len(all_docs):
        # Nothing removed
        return 0
    
    # Clear the entire index first
    clear_session_material(course_id, division)
    
    # If there are docs left, rebuild cache
    if remaining_docs:
        new_vs = FAISS.from_documents(remaining_docs, embeddings)
        key = f"{course_id}_{division}"
        faiss_indexes[key] = new_vs
        new_vs.save_local(str(get_index_path(course_id, division)))
    
    # 3. Also delete original file if it exists in materials
    file_path = config.MATERIALS_BASE_PATH / f"{course_id}_{division}" / source_name
    if file_path.exists():
        try:
            file_path.unlink()
        except Exception as e:
            print(f"DEBUG: Failed to delete local file {file_path}: {e}")
    
    return len(all_docs) - len(remaining_docs)

def clear_session_material(course_id: str, division: str):
    """Removes the entire FAISS index from memory and disk for a specific course and division."""
    key = f"{course_id}_{division}"
    
    # Clear from memory
    if key in faiss_indexes:
        del faiss_indexes[key]
    
    # Clear FAISS index from disk
    index_path = get_index_path(course_id, division)
    if os.path.exists(index_path):
        import shutil
        shutil.rmtree(index_path, ignore_errors=True)
    
    # Clear original materials from disk
    mat_path = config.MATERIALS_BASE_PATH / key
    if mat_path.exists():
        try:
            shutil.rmtree(mat_path)
        except Exception as e:
            print(f"DEBUG: Failed to remove materials directory {mat_path}: {e}")

def get_session_materials(course_id: str, division: str):
    """Returns a list of files and URLs currently in the session for UI sync."""
    key = f"{course_id}_{division}"
    files = []
    urls = []
    
    # Files from disk
    mat_path = config.MATERIALS_BASE_PATH / key
    if mat_path.exists():
        files = [f.name for f in mat_path.iterdir() if f.is_file()]
        
    # URLs from metadata in vector store (if loaded)
    vector_store = get_vector_store(course_id, division)
    if vector_store:
        all_docs = list(vector_store.docstore._dict.values())
        unique_urls = set()
        for doc in all_docs:
            src = doc.metadata.get('source', '')
            if src.startswith('http'):
                unique_urls.add(src)
        urls = list(unique_urls)
        
    return {"files": files, "urls": urls}

def clear_all_materials():
    """Purges ALL faculty sessions and materials. Called on logout."""
    # Clear all in-memory FAISS indices
    faiss_indexes.clear()
    
    # Clear all directories in materials base path
    if config.MATERIALS_BASE_PATH.exists():
        for item in config.MATERIALS_BASE_PATH.iterdir():
            if item.is_dir():
                try:
                    shutil.rmtree(item)
                except Exception as e:
                    print(f"DEBUG: Failed to remove directory {item} during clear_all: {e}")
